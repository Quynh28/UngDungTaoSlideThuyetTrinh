import base64
import glob
import os
import random
import string
import requests
import streamlit as st
from urllib.parse import urlparse
from PIL import Image
import openai
from icrawler import ImageDownloader
from icrawler.builtin import GoogleImageCrawler
from pptx import Presentation
from io import BytesIO
import datetime

# Cài đặt tông màu và giao diện chính
st.set_page_config(page_title="Tạo slide PPT tự động", page_icon="🔧", layout="wide")

# Áp dụng CSS tùy chỉnh
custom_css = """
<style>
.css-18e3th9 { background-color: #f8f9fa; }
.css-1d391kg { background-color: #cce9ff; color: black; width: 500px !important; }
.css-hxt7ib { font-size: 2rem; color: #343a40; font-weight: bold; }
.stButton>button { background-color: #343a40; color: white; border-radius: 5px; border: 1px solid #ced4da; }
.stTextInput>div>input { border-radius: 5px; border: 1px solid #d1d1d1; padding: 10px; }
</style>
"""
st.markdown(custom_css, unsafe_allow_html=True)

# Thanh tiêu đề
st.title("Tạo PPT bằng ChatGPT")

# Đọc danh sách các mẫu PowerPoint từ thư mục ppt_themes
def load_themes(directory):
    return [f for f in os.listdir(directory) if f.endswith('.pptx')]

ppt_themes_directory = "ppt_themes"
themes = load_themes(ppt_themes_directory)

# Thiết lập sidebar
st.sidebar.header("Chọn Mẫu:")
selected_theme = st.sidebar.selectbox("Chọn theme:", themes)

# Đường dẫn đến hình ảnh theme
image_folder = "theme_images"
image_path = os.path.join(image_folder, f"{os.path.splitext(selected_theme)[0]}.png")

# Kiểm tra và hiển thị hình ảnh theme
if os.path.exists(image_path):
    theme_image = Image.open(image_path)
    st.sidebar.image(theme_image, caption="Hình ảnh của theme", use_container_width=True)
else:
    st.sidebar.warning("Hình ảnh không tìm thấy!")

# Nội dung chính: Khu vực nhập liệu
st.header("Nhập thông tin")
title = st.text_input("Chủ đề thuyết trình:")
slide_length = st.number_input("Số lượng trang:", min_value=1, step=1, format="%d")

# Tạo cột để nhập link hình ảnh và tải lên hình ảnh
col1, col2 = st.columns(2)
with col1:
    image_urls = st.text_area("Nhập link hình ảnh (cách nhau bằng dấu phẩy):", placeholder="https://example.com/image.png")
with col2:
    uploaded_images = st.file_uploader("Tải lên hình ảnh (có thể chọn nhiều)", type=["jpg", "jpeg", "png"],
                                       accept_multiple_files=True)

# Hàm tạo tên ảnh duy nhất
def generate_unique_image_name():
    return ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in range(16))

class PrefixNameDownloader(ImageDownloader):
    def __init__(self, *args, **kwargs):
        self.image_name_prefix = generate_unique_image_name()
        super().__init__(*args, **kwargs)

    def get_filename(self, task, default_ext):
        url_path = urlparse(task['file_url'])[2]
        if '.' in url_path:
            extension = url_path.split('.')[-1]
            if extension.lower() not in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'ppm', 'pgm']:
                extension = default_ext
        else:
            extension = default_ext

        filename = base64.b64encode(url_path.encode()).decode()
        return f"p_{self.image_name_prefix}{{}}.{extension}"

# Biến lưu trữ buffer PPTX
pptx_buffer = None

# Hàm tạo PPT với nội dung chất lượng cao
def generate_ppt(topic, slide_length, selected_theme, uploaded_images, image_urls):
    # Tạo thư mục mới cho slide
    timestamp = datetime.datetime.now().strftime("%d%m%Y_%H%M%S")
    new_directory = f"slide_creation_{timestamp}"
    os.makedirs(new_directory, exist_ok=True)  # Tạo thư mục mới

    root = Presentation(os.path.join(ppt_themes_directory, selected_theme))
    openai.api_key = 'API_KEY'  # Thay thế bằng API key thực tế của bạn

    message = f"""Tạo một dàn bài cho một bài thuyết trình slideshow về chủ đề {topic} gồm {slide_length} trang. 
    Đảm bảo rằng mỗi trang có khoảng 100 từ và sử dụng tiếng Việt Nam.
    Bạn được phép sử dụng các loại slide sau:
    Slide tiêu đề - (Tiêu đề, Phụ đề)
    Slide nội dung - (Tiêu đề, Nội dung)
    Slide hình ảnh - (Tiêu đề, Nội dung, Hình ảnh)
    Slide cảm ơn - (Tiêu đề)

    Đặt thẻ này trước Slide tiêu đề: [L_TS]
    Đặt thẻ này trước Slide nội dung: [L_CS]
    Đặt thẻ này trước Slide hình ảnh: [L_IS]
    Đặt thẻ này trước Slide cảm ơn: [L_THS]

    Đặt thẻ này trước Tiêu đề: [TITLE]
    Đặt thẻ này sau Tiêu đề: [/TITLE]

    Đặt thẻ này trước Phụ đề: [SUBTITLE]
    Đặt thẻ này sau Phụ đề: [/SUBTITLE]

    Đặt thẻ này trước Nội dung: [CONTENT]
    Đặt thẻ này sau Nội dung: [/CONTENT]

    Đặt thẻ này trước Hình ảnh: [IMAGE]
    Đặt thẻ này sau Hình ảnh: [/IMAGE]

    Đặt "[SLIDEBREAK]" sau mỗi slide"""

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": message}],
    )

    def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    def set_font(placeholder, font_name='Times New Roman'):
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name

    def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        set_font(title_shape, 'Times New Roman')

        if slide.placeholders:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            set_font(subtitle_shape, 'Times New Roman')

    def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        set_font(slide.shapes.title, 'Times New Roman')

    def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        set_font(title_shape, "Times New Roman")

        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            content_shape.text = content
            set_font(content_shape, "Times New Roman")

    def create_title_and_content_and_image_slide(title, content, image):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        set_font(title_shape, 'Times New Roman')

        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            content_shape.text = content
            set_font(content_shape, "Times New Roman")

        # Kiểm tra xem hình ảnh là URL hay tệp tải lên
        if isinstance(image, str) and image.startswith("http"):
            response = requests.get(image)
            if response.status_code == 200:
                image_name = generate_unique_image_name() + ".jpg"
                image_path = os.path.join(new_directory, image_name)  # Lưu vào thư mục mới
                with open(image_path, 'wb') as f:
                    f.write(response.content)
                if len(slide.placeholders) > 2:
                    slide.shapes.add_picture(image_path,
                                             slide.placeholders[2].left,
                                             slide.placeholders[2].top,
                                             slide.placeholders[2].width,
                                             slide.placeholders[2].height)
            else:
                st.warning("Không thể tải hình ảnh từ URL.")
                return

        else:
            image_name = generate_unique_image_name() + ".png"
            image_path = os.path.join(new_directory, image_name)  # Lưu vào thư mục mới
            with Image.open(image) as img:
                img.save(image_path)

            if len(slide.placeholders) > 2:
                slide.shapes.add_picture(image_path,
                                         slide.placeholders[2].left,
                                         slide.placeholders[2].top,
                                         slide.placeholders[2].width,
                                         slide.placeholders[2].height)

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []

        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)

        return "".join(result) if result else ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def limit_words(content, max_words):
        words = content.split()
        if len(words) > max_words:
            truncated_content = ' '.join(words[:max_words])
            last_space = truncated_content.rfind(' ')
            if last_space != -1:
                truncated_content = truncated_content[:last_space]  # Cắt đến từ cuối cùng
            return truncated_content
        return content

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        image_index = 0
        downloaded_images = []  # Danh sách lưu trữ hình ảnh đã tải về
        image_urls_list = [url.strip() for url in image_urls.split(',')] if image_urls else []

        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)

            if slide_type == "[L_TS]":
                create_title_slide(
                    find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"),
                    find_text_in_between_tags(slide, "[SUBTITLE]", "[/SUBTITLE]")
                )

            elif slide_type == "[L_CS]":
                content = find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]")
                limited_content = limit_words(content, 100)  # Giới hạn 100 từ
                create_title_and_content_slide(
                    find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"),
                    limited_content
                )

            elif slide_type == "[L_IS]":
                title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")  # Lấy tiêu đề của slide
                content = find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]")
                image_to_use = None

                # Sử dụng hình ảnh từ URL hoặc tệp tải lên
                if image_urls_list:
                    image_to_use = image_urls_list[min(image_index, len(image_urls_list) - 1)]
                    image_index += 1

                elif uploaded_images:
                    image_to_use = uploaded_images[min(len(uploaded_images) - 1, len(downloaded_images))]  # Lấy hình ảnh tiếp theo

                if not image_to_use and uploaded_images:
                    image_to_use = uploaded_images[min(len(uploaded_images) - 1, len(downloaded_images))]

                if image_to_use:
                    create_title_and_content_and_image_slide(
                        title,
                        content,
                        image_to_use
                    )
                else:
                    enhanced_query = f"{title} - {content} - hình ảnh"
                    google_crawler = GoogleImageCrawler(downloader_cls=PrefixNameDownloader,
                                                        storage={'root_dir': new_directory})  # Lưu vào thư mục mới
                    google_crawler.crawl(keyword=enhanced_query, max_num=10)

                    dir_path = os.path.dirname(os.path.realpath(__file__))
                    file_names = glob.glob(os.path.join(new_directory, "p_*"))  # Lấy tất cả hình ảnh đã tải về

                    if not file_names:
                        st.error("Không tìm thấy hình ảnh cho truy vấn này.")
                        continue

                    downloaded_images.extend(file_names)  # Thêm tất cả hình ảnh vào danh sách

                    if image_index < len(downloaded_images):
                        image_to_use = downloaded_images[image_index]
                        image_index += 1

                    create_title_and_content_and_image_slide(
                        title,
                        content,
                        image_to_use
                    )

            elif slide_type == "[L_THS]":
                create_section_header_slide(find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]"))

    delete_all_slides()
    parse_response(response.choices[0].message.content)

    pptx_buffer = BytesIO()
    root.save(pptx_buffer)
    pptx_buffer.seek(0)

    return pptx_buffer, title

# Kiểm tra và hiển thị thông báo lỗi chỉ khi người dùng nhấn nút "Tạo PPT"
if st.button("Tạo PPT"):
    if title.strip() and slide_length > 0 and selected_theme:
        pptx_buffer, message = generate_ppt(title, slide_length, selected_theme, uploaded_images, image_urls)
        st.success("Đã xong! File PowerPoint của bạn đã được tạo.")

        # Thêm nút tải xuống
        st.download_button(
            label="Tải xuống PPT",
            data=pptx_buffer,
            file_name=f"{title}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    else:
        if not title.strip():
            st.error("Vui lòng nhập chủ đề thuyết trình.")
        if slide_length <= 0:
            st.error("Vui lòng nhập số lượng trang lớn hơn 0.")

# Footer
st.markdown("""<hr style='border: 1px solid #ced4da;'/>""", unsafe_allow_html=True)
st.caption("Nguyễn Hà Nghi Phương B2013553.<br>Nguyễn Thị Ngọc Quỳnh B2013556.", unsafe_allow_html=True)